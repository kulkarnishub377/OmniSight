from pathlib import Path
from datetime import date
import textwrap

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'reports'
OUT.mkdir(exist_ok=True)
TODAY = '2026-06-15'

report_md = """# AMIF v1.0 - Autonomous Multi-Modal Intelligence Fabric

## Complete Project Report

**Date:** {TODAY}  
**Project Type:** Senior-level full-stack AI system design and implementation project  
**Primary Domain:** Industrial safety, incident intelligence, predictive maintenance, and operational monitoring

---

## 1. Executive Summary

AMIF, or Autonomous Multi-Modal Intelligence Fabric, is a production-style AI incident intelligence platform. It ingests operational signals from cameras, sensors, audio streams, documents, and manual/API sources, normalizes those signals into canonical events, correlates them into incidents, retrieves relevant evidence from memory, runs an agentic investigation workflow, applies safety guardrails, and presents the result through a professional operator dashboard.

The project is designed as a resume-worthy senior full-stack and AI architecture project. The current implementation is a runnable MVP with strong architecture boundaries. Some AI services are implemented as deterministic/mock boundaries so the system can run locally without heavyweight models, but the architecture is ready to plug in real models such as YOLO, Whisper, BGE embeddings, Qdrant, LangGraph, Llama/Qwen/DeepSeek, and Llama Guard.

Core product loop:

```text
Detect -> Normalize -> Correlate -> Retrieve -> Reason -> Guard -> Alert -> Explain
```

---

## 2. Problem Statement

Industrial and enterprise environments produce many isolated signals: CCTV frames, IoT sensor readings, machine logs, audio warnings, maintenance manuals, emails, and operator notes. Traditional systems usually treat these streams separately. This creates problems:

- Operators receive too many isolated alerts.
- Root-cause analysis is slow and manual.
- Evidence is scattered across cameras, sensors, manuals, and incident history.
- AI decisions are often not auditable.
- Automated actions can be unsafe without policy guardrails.

AMIF solves this by acting as an incident intelligence fabric: it fuses multimodal signals, detects patterns, retrieves supporting knowledge, generates evidence-backed investigation summaries, and routes safe actions to humans or external systems.

---

## 3. Real-World Use Cases

### 3.1 Factory Safety Monitoring

Detect unsafe events such as forklift activity in restricted zones, workers entering danger areas, or machinery overheating. Correlate visual detections with sensor anomalies and create high-priority incidents.

### 3.2 Predictive Maintenance

Monitor temperature, vibration, noise, and energy readings. When abnormal readings occur, retrieve machine manuals and past incident reports to recommend inspection steps.

### 3.3 Warehouse Operations

Detect unsafe vehicle movement, zone breaches, PPE violations, and abnormal activity near equipment. Automatically alert supervisors and create tickets.

### 3.4 Security Operations

Correlate unauthorized access, camera detections, sensor events, and operator notes. Generate evidence-backed security incident reports.

### 3.5 Compliance and Audit

Store all events, alerts, actions, agent runs, guard decisions, and operator acknowledgements. Provide a traceable audit trail for regulated environments.

### 3.6 Smart Facilities

Combine cameras, HVAC sensors, fire alarms, access systems, and maintenance documents to detect and investigate facility incidents.

---

## 4. High-Level Architecture

```text
Data Sources
  -> Connectors / Ingestion APIs
  -> FastAPI API Gateway
  -> Domain Services
       - Event Service
       - Vision Service
       - Audio Service
       - Sensor Service
       - Document Service
       - Search Service
       - Agent Service
       - Notification Service
       - Action Service
       - User/Auth Service
  -> Redpanda-ready Event Bus Topics
  -> Event Processing Layer
       - validation
       - deduplication
       - enrichment
       - anomaly generation
       - time-window correlation
  -> Memory Fabric
       - PostgreSQL historical memory
       - Redis episodic memory design
       - Qdrant semantic memory design
       - Neo4j-style knowledge graph endpoint
  -> Agent Runtime
       - Observer
       - Retriever
       - Investigator
       - Planner
       - Safety Guard
       - Executor
  -> Operator Console / Dashboard
```

---

## 5. Technology Stack

### Backend

- Python
- FastAPI
- Pydantic
- SQLAlchemy
- JWT authentication
- RBAC middleware pattern

### Storage

- PostgreSQL-compatible SQLAlchemy models
- SQLite fallback for local fast development
- Redis-ready episodic memory design
- Qdrant-ready semantic memory design
- Neo4j-style graph API

### Eventing

- Redpanda included in Docker Compose
- Redpanda-ready topic contracts
- In-process event processing for MVP
- FastStream extraction path documented

### AI / Intelligence Layer

- YOLO/RT-DETR-style vision service boundary
- Whisper-style audio service boundary
- OCR/document intelligence pipeline boundary
- Deterministic embedding-style semantic retrieval fallback
- Agent runtime with Observer, Retriever, Investigator, Planner, Guard, Executor
- Llama Guard-style safety policy boundary

### Frontend

- Advanced vanilla JavaScript SPA
- CSS design system
- Responsive dashboard
- Animated KPI cards
- Command palette
- Interactive SVG knowledge graph
- Modals and toast notifications

### Observability / Deployment

- Prometheus metrics endpoint
- Grafana and Prometheus included in Compose
- Docker Compose
- Kubernetes starter manifests

---

## 6. Implemented Project Structure

```text
amif/
  backend/
    app/
      main.py
      core/
        config.py
        security.py
      db/
        session.py
        init_db.py
      models/
        models.py
        enums.py
      schemas/
        schemas.py
      routers/
        auth.py
        events.py
        incidents.py
        alerts.py
        documents.py
        audio.py
        actions.py
        audit.py
        users.py
        notifications.py
        knowledge.py
        system.py
        demo.py
        health.py
      services/
        event_service.py
        event_processor.py
        document_service.py
        search_service.py
        agent_service.py
        action_service.py
        audit_service.py
      static/
        index.html
        architecture.html
        assets/
          app.css
          app.js
  connectors/
    camera_connector.py
    audio_connector.py
    iot_connector.py
    document_connector.py
    email_connector.py
  workers/
    event_processor/
      README.md
  infra/
    k8s/
    prometheus/
  docs/
  reports/
```

---

## 7. Core Backend Modules

### 7.1 API Gateway

The FastAPI app acts as the gateway for the current MVP. It provides authentication, routing, OpenAPI documentation, static frontend hosting, health checks, metrics, and API endpoints.

Important endpoints:

```text
POST /api/auth/login
GET  /api/auth/me
POST /api/events
GET  /api/events
POST /api/vision/mock-detect
POST /api/audio/mock-transcribe
POST /api/sensors/temperature
GET  /api/incidents
PATCH /api/incidents/{incident_id}
POST /api/incidents/{incident_id}/investigate
GET  /api/incidents/{incident_id}/agent-runs
GET  /api/alerts
POST /api/alerts/{alert_id}/acknowledge
POST /api/documents/upload
POST /api/search/query
GET  /api/knowledge/graph
GET  /api/system/stats
GET  /api/audit/logs
```

### 7.2 Event Service

The Event Service is responsible for accepting canonical AMIF events, validating inputs, storing event metadata, logging audits, and triggering event processing.

The canonical AMIF event includes:

```json
{
  "schema_version": "1.0",
  "event_id": "evt_xxx",
  "source_id": "camera_warehouse_a_01",
  "source_type": "camera",
  "source_sequence": 123,
  "event_type": "forklift_detected",
  "severity": "medium",
  "timestamp": "2026-06-01T10:30:00Z",
  "tenant_id": "demo_tenant",
  "location": {"site": "plant_1", "zone": "restricted_zone"},
  "payload": {"object": "forklift", "confidence": 0.94},
  "trace": {"producer": "vision-service"}
}
```

### 7.3 Event Processor

The processor handles:

- event status updates
- temperature anomaly derivation
- time-window correlation
- incident creation
- alert creation
- ticket stub creation
- agent investigation triggering

Implemented correlation rule:

```text
IF forklift_detected in restricted_zone
AND temperature_anomaly occurs within 10 minutes
THEN create high-severity incident
```

### 7.4 Document Service

The Document Service supports text document upload, cleanup, chunking, metadata creation, and indexing into the semantic retrieval fallback.

### 7.5 Search Service

The Search Service implements an offline-friendly deterministic semantic retrieval system using hashing vectors and lexical overlap. It is designed to be replaced or extended with Qdrant and real embeddings.

### 7.6 Agent Service

The Agent Service implements a deterministic agent workflow:

```text
Observer -> Retriever -> Investigator -> Planner -> SafetyGuard -> ActionExecutor
```

The agent generates:

- incident summary
- hypotheses
- evidence references
- risk score
- recommended actions
- guard decision
- trace steps

### 7.7 Action Service

The Action Service currently supports:

- dashboard notification stub
- ticket creation stub
- action result storage
- audit logging

It can later be extended to Jira, ServiceNow, Slack, Teams, SAP, webhooks, and email.

---

## 8. Frontend Operator Console

The AMIF frontend is a professional animated vanilla-JS SPA served by FastAPI.

### Implemented Screens

1. Login
2. Overview dashboard
3. Incident command center
4. Event operations
5. Documents and RAG
6. Agent runtime
7. Knowledge graph
8. Observability
9. Security and audit

### Advanced UI Features

- responsive sidebar navigation
- animated KPI counters
- command palette with Ctrl/Cmd + K
- light/dark theme toggle
- live refresh toggle
- workspace JSON export
- modal detail inspectors
- event detail modal
- graph node modal
- toast notifications
- canvas sparklines
- interactive SVG knowledge graph
- confetti animation after demo seed
- progress loading strip
- skeleton loading states
- global search

---

## 9. AI Used in the Project

Important clarification: the current MVP uses AI-ready service boundaries and deterministic mocks/fallbacks so it can run locally. It does not require heavyweight model downloads.

### 9.1 Vision AI

- Boundary: YOLO/RT-DETR-style object detection
- Current: mock forklift/person detection event generation
- Real integration path: YOLOv8/YOLOv11, RT-DETR, Florence-2

### 9.2 Audio AI

- Boundary: Whisper-style speech-to-text
- Current: mock transcription event generation
- Real integration path: Whisper, faster-whisper

### 9.3 Document AI

- Boundary: OCR, chunking, embedding generation
- Current: text upload, chunking, deterministic vector fallback
- Real integration path: PyMuPDF, PaddleOCR/Tesseract, BGE-M3 embeddings

### 9.4 RAG / Semantic Retrieval

- Current: local hashing vector search and keyword overlap
- Real integration path: Qdrant vector database, reranking, metadata filtering

### 9.5 Agentic AI

- Boundary: LangGraph-style agent workflow
- Current: deterministic agents for reliable local demos
- Real integration path: LangGraph plus Llama/Qwen/DeepSeek models

### 9.6 Safety AI

- Boundary: Llama Guard-style safety/policy validation
- Current: deterministic guard checks requiring human approval for risky actions
- Real integration path: Llama Guard or enterprise policy classifier

---

## 10. Data Flow Demo Scenario

The demo scenario works as follows:

1. User logs in as Admin.
2. User clicks Seed Demo.
3. System uploads Machine A safety manual.
4. System emits a forklift detection event from camera.
5. System emits a temperature reading from Machine A.
6. Processor derives a temperature anomaly.
7. Processor correlates forklift in restricted zone with overheating machine.
8. System creates a high-severity incident.
9. System creates an alert and ticket stub.
10. Agent runtime investigates the incident.
11. Search service retrieves relevant manual chunks.
12. Agent creates hypotheses, risk score, and action plan.
13. Safety Guard blocks actions requiring human approval.
14. Dashboard displays summary, evidence, graph, alerts, and audit trail.

---

## 11. Memory Fabric

### PostgreSQL / SQLite

Used as source of truth for:

- users
- events
- incidents
- incident-event relationships
- alerts
- documents
- document chunks
- agent runs
- action results
- audit logs

### Redis-ready Episodic Memory

Designed for:

- recent events
- active incidents
- agent scratchpad state
- dedupe keys
- live dashboard cache

### Qdrant-ready Semantic Memory

Designed for:

- manual chunks
- safety policies
- incident summaries
- maintenance reports
- similarity search

### Knowledge Graph Endpoint

The `/api/knowledge/graph` endpoint exposes Neo4j-style relationships:

- Incident HAS_EVIDENCE Event
- Event AFFECTS Machine
- Event LOCATED_IN Zone
- Document DESCRIBES Machine

---

## 12. Security and Governance

Implemented controls:

- JWT authentication
- password hashing
- role-based access control
- Admin, Operator, Analyst, Viewer roles
- audit logging
- action traceability
- incident update audit trail

Role capabilities:

```text
Admin: full access, users, audit logs
Operator: acknowledge alerts, update incidents
Analyst: upload documents, trigger investigations
Viewer: read-only dashboard access
```

---

## 13. Observability

Implemented:

- `/health` endpoint
- `/metrics` Prometheus endpoint
- system stats endpoint
- Prometheus config
- Grafana service in Docker Compose
- dashboard observability screen

Tracked platform indicators:

- event count
- incident count
- alert count
- document count
- agent run count
- audit log count

Recommended future metrics:

- p50/p95/p99 event latency
- consumer lag
- DLQ rate
- agent latency
- RAG retrieval precision
- false alert rate
- guard rejection rate

---

## 14. Deployment

### Local Run

```bash
cd amif
cp .env.example .env
docker compose up --build
```

Open:

```text
Dashboard:     http://localhost:8000
Architecture:  http://localhost:8000/architecture.html
API Docs:      http://localhost:8000/docs
Prometheus:    http://localhost:9090
Grafana:       http://localhost:3000
```

### Kubernetes Readiness

Starter manifests are included for:

- namespace
- backend deployment
- backend service
- horizontal pod autoscaler

Future Kubernetes work:

- split services into separate deployments
- add ingress
- add secrets
- add network policies
- add persistent volumes
- add Helm chart

---

## 15. Current Limitations

The MVP is intentionally optimized for local portfolio demonstration.

Current limitations:

- event processing is in-process, not yet a real Redpanda consumer
- vision/audio models are mocked boundaries
- embeddings are deterministic local fallback, not real neural embeddings
- Qdrant and Redis are infrastructure-ready but not deeply integrated
- no real external Jira/Slack/ServiceNow integration yet
- no production-grade migration system yet
- no full test suite yet

---

## 16. Future Roadmap

### Phase 1: Real Streaming

- Add FastStream worker
- Consume from Redpanda topics
- Add DLQ processing
- Add replay tooling

### Phase 2: Real AI Models

- Add YOLO inference service
- Add Whisper/faster-whisper service
- Add BGE-M3 embeddings
- Add Qdrant indexing

### Phase 3: LangGraph Agents

- Replace deterministic agent workflow with LangGraph
- Add model gateway
- Add prompt versioning
- Add model fallback

### Phase 4: Enterprise Integrations

- Jira ticket creation
- Slack/Teams alerts
- ServiceNow integration
- email notifications
- SAP/CMMS hooks

### Phase 5: Production Hardening

- Alembic migrations
- test suite
- OpenTelemetry distributed traces
- Kubernetes Helm chart
- load testing
- security hardening

---

## 17. Resume Positioning

Suggested resume project title:

**Autonomous Multi-Modal Incident Intelligence Fabric**

Suggested bullets:

- Architected and implemented an event-driven AI incident intelligence platform using FastAPI, SQLAlchemy, Redpanda-ready topic contracts, PostgreSQL, Redis/Qdrant-ready memory layers, and an advanced vanilla-JS operator console.
- Built canonical event ingestion, validation, anomaly derivation, time-window correlation, incident creation, alerting, audit logging, and agent investigation workflows.
- Designed a multi-layer memory fabric combining operational storage, semantic retrieval, episodic memory design, and a knowledge graph API.
- Implemented an auditable agent runtime with Observer, Retriever, Investigator, Planner, Safety Guard, and Executor stages.
- Developed a production-style frontend with command palette, animated KPIs, event explorer, RAG interface, agent trace viewer, interactive graph visualization, observability, and RBAC/audit screens.
- Added Docker Compose infrastructure with Postgres, Redis, Qdrant, Redpanda, Prometheus, and Grafana plus Kubernetes starter manifests.

---

## 18. Final Summary

AMIF is a complete full-stack AI architecture project that demonstrates senior-level system design, backend engineering, frontend UX, AI orchestration, event processing, memory design, governance, and deployment readiness.

The project is not just a dashboard. It is an end-to-end incident intelligence platform with clear contracts, extensible service boundaries, auditable agent workflows, and a polished operator experience.
"""
report_md = report_md.replace('{TODAY}', TODAY)

(OUT / 'AMIF_Complete_Project_Report.md').write_text(report_md, encoding='utf-8')

# HTML report
html = f"""<!doctype html>
<html><head><meta charset='utf-8'><title>AMIF Complete Project Report</title>
<style>
body{{font-family:Inter,Arial,sans-serif;background:#f7f9fc;color:#172033;line-height:1.55;margin:0}}
.cover{{background:linear-gradient(135deg,#071126,#173a76);color:white;padding:56px 70px}}
.cover h1{{font-size:42px;margin:0 0 8px}}.cover p{{opacity:.88;font-size:18px}}
main{{max-width:980px;margin:auto;background:white;padding:42px 58px;box-shadow:0 12px 50px #10204018}}
h1,h2,h3{{color:#0d2350}} h2{{border-top:1px solid #d8e0ef;padding-top:24px;margin-top:34px}}
pre,code{{background:#0b1222;color:#dbeafe;border-radius:10px;padding:12px;display:block;overflow:auto}}
blockquote{{border-left:4px solid #38bdf8;background:#eff8ff;margin:18px 0;padding:12px 16px}}
li{{margin:4px 0}} table{{border-collapse:collapse;width:100%}}td,th{{border:1px solid #d8e0ef;padding:8px}}
.badge{{display:inline-block;padding:4px 8px;border-radius:999px;background:#e0f2fe;color:#075985;font-weight:700;font-size:12px}}
</style></head><body><div class='cover'><h1>AMIF v1.0</h1><p>Autonomous Multi-Modal Intelligence Fabric</p><p>Complete Project Report - {TODAY}</p></div><main>
"""
# basic markdown -> html conversion for report without external markdown lib
lines = report_md.splitlines()
in_code = False
for line in lines:
    if line.startswith('```'):
        html += '</pre>' if in_code else '<pre>'
        in_code = not in_code
        continue
    if in_code:
        html += line.replace('&','&amp;').replace('<','&lt;').replace('>','&gt;') + '\n'
        continue
    if line.startswith('# '): html += f"<h1>{line[2:]}</h1>\n"
    elif line.startswith('## '): html += f"<h2>{line[3:]}</h2>\n"
    elif line.startswith('### '): html += f"<h3>{line[4:]}</h3>\n"
    elif line.startswith('- '): html += f"<ul><li>{line[2:]}</li></ul>\n"
    elif line.strip() == '---': html += '<hr>\n'
    elif line.strip(): html += f"<p>{line}</p>\n"
html += '</main></body></html>'
(OUT / 'AMIF_Complete_Project_Report.html').write_text(html, encoding='utf-8')

# PDF generation with ReportLab
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle, Preformatted, KeepTogether
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfbase import pdfmetrics

pdf_path = OUT / 'AMIF_Complete_Project_Report.pdf'
doc = SimpleDocTemplate(str(pdf_path), pagesize=A4, rightMargin=42, leftMargin=42, topMargin=48, bottomMargin=42)
styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name='CoverTitle', parent=styles['Title'], fontSize=28, leading=34, textColor=colors.HexColor('#0F2D5C'), spaceAfter=12))
styles.add(ParagraphStyle(name='Section', parent=styles['Heading2'], fontSize=15, leading=19, textColor=colors.HexColor('#123B73'), spaceBefore=14, spaceAfter=8))
styles.add(ParagraphStyle(name='SubSection', parent=styles['Heading3'], fontSize=12, leading=15, textColor=colors.HexColor('#1E4C86'), spaceBefore=8, spaceAfter=5))
styles.add(ParagraphStyle(name='BodyX', parent=styles['BodyText'], fontSize=9.4, leading=13.2, spaceAfter=6))
styles.add(ParagraphStyle(name='BulletX', parent=styles['BodyText'], fontSize=9.2, leading=12.4, leftIndent=14, firstLineIndent=-8, spaceAfter=3))
styles.add(ParagraphStyle(name='Small', parent=styles['BodyText'], fontSize=8, leading=10, textColor=colors.HexColor('#4B5563')))

story = []
story.append(Spacer(1, 0.7*inch))
story.append(Paragraph('AMIF v1.0', styles['CoverTitle']))
story.append(Paragraph('Autonomous Multi-Modal Intelligence Fabric', styles['Heading2']))
story.append(Paragraph('Complete Project Report', styles['Heading3']))
story.append(Spacer(1, 0.3*inch))
story.append(Paragraph(f'Date: {TODAY}', styles['BodyX']))
story.append(Paragraph('Project Type: Senior-level full-stack AI system design and implementation project', styles['BodyX']))
story.append(Paragraph('Primary Domain: Industrial safety, incident intelligence, predictive maintenance, and operational monitoring', styles['BodyX']))
story.append(Spacer(1, 0.4*inch))
story.append(Table([[Paragraph('Core Loop', styles['Heading3']), Paragraph('Detect -> Normalize -> Correlate -> Retrieve -> Reason -> Guard -> Alert -> Explain', styles['BodyX'])]], colWidths=[1.3*inch, 4.9*inch], style=TableStyle([('BACKGROUND',(0,0),(-1,-1),colors.HexColor('#EAF5FF')),('BOX',(0,0),(-1,-1),0.5,colors.HexColor('#94A3B8')),('VALIGN',(0,0),(-1,-1),'TOP'),('PADDING',(0,0),(-1,-1),10)])))
story.append(PageBreak())

sections = [
('1. Executive Summary', ['AMIF is a production-style AI incident intelligence platform. It ingests operational signals from cameras, sensors, audio streams, documents, and APIs; normalizes those signals into events; correlates events into incidents; retrieves evidence; runs agentic investigation; applies guardrails; and presents results through an advanced operator console.', 'The current implementation is a runnable MVP with AI-ready service boundaries and deterministic local fallbacks so it works without heavyweight models.']),
('2. Business Problem', ['Industrial teams receive fragmented alerts from cameras, sensors, logs, and documents. AMIF reduces alert fatigue and speeds root-cause analysis by fusing multimodal evidence into auditable incident intelligence.']),
('3. Real Use Cases', ['Factory safety monitoring: detect forklifts or people in restricted zones.', 'Predictive maintenance: detect overheating, vibration, or abnormal noise and retrieve manuals.', 'Warehouse operations: monitor unsafe vehicle movement, PPE gaps, and zone breaches.', 'Security operations: correlate access, camera, and sensor evidence.', 'Compliance: maintain audit trails of events, alerts, actions, and agent decisions.', 'Smart facilities: combine HVAC, alarms, cameras, access systems, and documents.']),
('4. Architecture Overview', ['Data Sources -> Connectors -> FastAPI Gateway -> Domain Services -> Redpanda-ready Event Bus -> Event Processing -> Memory Fabric -> Agent Runtime -> Operator Console.', 'The MVP runs as a modular FastAPI application while preserving clean seams for future service extraction.']),
('5. Backend Services', ['API Gateway/Auth: login, JWT, RBAC, docs, metrics.', 'Event Service: canonical AMIF event ingestion and persistence.', 'Event Processor: validation, anomaly derivation, correlation, incident generation.', 'Document Service: upload, cleanup, chunking, indexing.', 'Search Service: deterministic semantic retrieval fallback.', 'Agent Service: Observer, Retriever, Investigator, Planner, Guard, Executor.', 'Action Service: ticket stub, dashboard notification, action audit.', 'Knowledge Service: Neo4j-style graph endpoint.']),
('6. Frontend Console', ['Advanced vanilla-JS SPA with login, overview, incident command, event operations, documents/RAG, agent runtime, knowledge graph, observability, and security/audit screens.', 'UX features include command palette, animated KPI counters, live refresh, theme toggle, modals, event explorer, SVG graph, sparkline charts, export, toasts, and responsive layout.']),
('7. AI Components', ['Vision AI boundary: YOLO/RT-DETR style detection; current demo mocks forklift detection.', 'Audio AI boundary: Whisper-style transcription; current demo mocks audio transcript events.', 'Document AI boundary: OCR/chunking/embeddings; current demo uses text chunking and vector fallback.', 'RAG: semantic retrieval over document chunks and incident evidence.', 'Agentic AI: LangGraph-style investigation workflow implemented deterministically.', 'Safety AI: Llama Guard-style policy checks and human approval logic.']),
('8. Demo Scenario', ['Seed Demo uploads a Machine A manual, emits a forklift detection, emits a temperature reading, derives a temperature anomaly, correlates both within a 10-minute window, creates a high-severity incident, opens an alert, creates a ticket stub, runs agent investigation, retrieves evidence, and displays summary, trace, graph, and audit logs.']),
('9. Data and Memory', ['PostgreSQL/SQLite stores users, events, incidents, alerts, documents, agent runs, actions, and audit logs.', 'Redis-ready design covers recent events, active incident context, and agent scratchpad.', 'Qdrant-ready design covers semantic memory over documents and incidents.', 'Knowledge graph endpoint exposes Incident-Event-Machine-Zone-Document relationships.']),
('10. Security and Governance', ['JWT authentication, hashed passwords, RBAC, audit logs, incident update audit trail, agent run trace storage, and guardrail checks for actions requiring human approval.']),
('11. Observability and Deployment', ['Health endpoint, Prometheus metrics endpoint, system stats API, Prometheus config, Grafana service, Docker Compose, and Kubernetes starter manifests are included.']),
('12. Limitations', ['AI models are mocked/deterministic in the MVP.', 'Event processing is currently in-process instead of a dedicated Redpanda consumer.', 'Redis and Qdrant are infra-ready but not fully integrated.', 'External integrations are stubs.', 'Test suite and migrations are future work.']),
('13. Roadmap', ['Add FastStream Redpanda worker and DLQ.', 'Integrate YOLO, Whisper, BGE embeddings, Qdrant.', 'Replace deterministic workflow with LangGraph and real LLMs.', 'Add Slack/Jira/ServiceNow/SAP integrations.', 'Add OpenTelemetry, Alembic, tests, Helm chart, and load testing.']),
('14. Resume Value', ['AMIF demonstrates senior-level full-stack engineering, system design, event-driven architecture, AI orchestration, RAG, agent workflows, safety governance, observability, deployment, and frontend UX.']),
]
for title, bullets in sections:
    story.append(Paragraph(title, styles['Section']))
    for b in bullets:
        story.append(Paragraph(f'- {b}', styles['BulletX']))

story.append(Paragraph('Key API Endpoints', styles['Section']))
api_rows = [
['Endpoint', 'Purpose'],
['POST /api/auth/login', 'Authenticate and receive JWT'],
['POST /api/demo/seed', 'Run complete demo scenario'],
['POST /api/events', 'Ingest canonical event'],
['GET /api/incidents', 'List incidents'],
['POST /api/incidents/{id}/investigate', 'Run agent investigation'],
['POST /api/documents/upload', 'Upload and index document'],
['POST /api/search/query', 'RAG semantic search'],
['GET /api/knowledge/graph', 'Graph relationships'],
['GET /api/system/stats', 'Dashboard system stats'],
]
story.append(Table(api_rows, colWidths=[2.35*inch, 3.85*inch], style=TableStyle([('BACKGROUND',(0,0),(-1,0),colors.HexColor('#123B73')),('TEXTCOLOR',(0,0),(-1,0),colors.white),('GRID',(0,0),(-1,-1),0.25,colors.HexColor('#CBD5E1')),('BACKGROUND',(0,1),(-1,-1),colors.HexColor('#F8FAFC')),('FONT',(0,0),(-1,-1),'Helvetica'),('FONTSIZE',(0,0),(-1,-1),8.4),('PADDING',(0,0),(-1,-1),6)])))
story.append(Spacer(1, 0.2*inch))
story.append(Paragraph('Final Summary', styles['Section']))
story.append(Paragraph('AMIF is an end-to-end incident intelligence platform with clear contracts, AI-ready services, auditable agent workflows, and a polished operator experience. It is suitable as a senior-level portfolio project and as a foundation for production-grade industrial AI systems.', styles['BodyX']))

def page_num(canvas, doc):
    canvas.saveState()
    canvas.setFont('Helvetica', 8)
    canvas.setFillColor(colors.HexColor('#64748B'))
    canvas.drawRightString(A4[0]-42, 24, f'AMIF Complete Project Report - Page {doc.page}')
    canvas.restoreState()

doc.build(story, onFirstPage=page_num, onLaterPages=page_num)

# PPTX generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ppt_path = OUT / 'AMIF_Project_Presentation.pptx'
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(7, 11, 22)
PANEL = RGBColor(17, 24, 43)
CYAN = RGBColor(103, 232, 249)
WHITE = RGBColor(238, 243, 255)
MUTED = RGBColor(148, 163, 184)
VIOLET = RGBColor(167, 139, 250)
GREEN = RGBColor(52, 211, 153)
RED = RGBColor(251, 113, 133)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = WHITE
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.35))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle; p2.font.size = Pt(11); p2.font.color.rgb = MUTED


def add_bullets(slide, bullets, x=0.75, y=1.45, w=11.8, h=5.6, font=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b; p.level = 0; p.font.size = Pt(font); p.font.color.rgb = WHITE if not b.startswith('  ') else MUTED
        p.space_after = Pt(8)


def add_card(slide, x, y, w, h, title, body, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = color
    tx = slide.shapes.add_textbox(Inches(x+0.15), Inches(y+0.12), Inches(w-0.3), Inches(h-0.2))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = color
    p2 = tf.add_paragraph(); p2.text = body; p2.font.size = Pt(10.5); p2.font.color.rgb = WHITE

slides = [
('AMIF v1.0', 'Autonomous Multi-Modal Intelligence Fabric', ['Senior-level full-stack AI architecture project', 'Industrial safety, incident intelligence, predictive maintenance', 'Core loop: Detect -> Correlate -> Retrieve -> Reason -> Guard -> Alert']),
('Problem', None, ['Operational signals are fragmented across cameras, sensors, audio, documents, and enterprise systems.', 'Operators receive isolated alerts without evidence-backed context.', 'Root-cause investigation is slow, manual, and hard to audit.', 'Autonomous actions require governance and human approval controls.']),
('Solution', None, ['AMIF turns multimodal signals into explainable incident intelligence.', 'It correlates camera, sensor, audio, document, and manual/API events.', 'It retrieves relevant manuals and prior context using semantic memory.', 'It runs an auditable agent workflow and produces safe action recommendations.']),
('High-Level Architecture', None, ['Data Sources -> Connectors -> API Gateway', 'Domain Services -> Redpanda-ready Event Bus', 'Event Processing -> Memory Fabric', 'Agent Runtime -> Safety Guard -> Action Engine', 'Operator Console -> Audit and Observability']),
('Implemented Backend', None, ['FastAPI gateway with JWT auth and RBAC', 'SQLAlchemy models for users, events, incidents, alerts, documents, agent runs, actions, audit logs', 'Event ingestion, processing, anomaly derivation, and incident correlation', 'Document chunking and deterministic semantic search fallback', 'Agent runtime and guardrail workflow']),
('Frontend Console', None, ['Advanced vanilla-JS single-page dashboard', 'Animated KPIs, command palette, live refresh, theme toggle', 'Incident command center, event explorer, RAG UI, agent trace viewer', 'Interactive SVG knowledge graph, observability, security, audit screens']),
('AI Components', None, ['Vision AI boundary: YOLO/RT-DETR style detection', 'Audio AI boundary: Whisper-style transcription', 'Document AI: OCR/chunking/embeddings pipeline boundary', 'RAG: semantic retrieval over manuals and incident evidence', 'Agentic AI: Observer, Retriever, Investigator, Planner, Guard, Executor', 'Safety AI: Llama Guard-style policy checks']),
('Demo Scenario', None, ['Upload Machine A safety manual', 'Emit forklift detection from restricted zone camera', 'Emit overheating sensor reading', 'Derive temperature anomaly and correlate events', 'Create high-severity incident and alert', 'Run agent investigation and display evidence-backed summary']),
('Memory Fabric', None, ['PostgreSQL/SQLite: durable source of truth', 'Redis-ready: recent events and active incident context', 'Qdrant-ready: semantic memory for manuals and past incidents', 'Neo4j-style endpoint: Incident-Event-Machine-Zone-Document graph']),
('Security and Governance', None, ['JWT authentication and password hashing', 'RBAC roles: Admin, Operator, Analyst, Viewer', 'Audit logs for events, actions, incidents, and agent runs', 'Guardrails block risky actions requiring human approval']),
('Observability and Deployment', None, ['Health and metrics endpoints', 'Prometheus and Grafana in Docker Compose', 'System stats API displayed in dashboard', 'Docker Compose local deployment', 'Kubernetes starter manifests included']),
('Current Limitations', None, ['AI services are mock/deterministic boundaries in MVP', 'Processing is in-process, not yet a Redpanda worker', 'Qdrant/Redis are infra-ready but not fully integrated', 'External integrations are action stubs', 'Migrations and full tests are future work']),
('Roadmap', None, ['Add FastStream + Redpanda worker and DLQ', 'Integrate YOLO, Whisper, BGE embeddings, Qdrant', 'Replace deterministic workflow with LangGraph and LLM gateway', 'Add Slack/Jira/ServiceNow/SAP integrations', 'Add OpenTelemetry, Alembic, tests, Helm, and load testing']),
('Resume Value', None, ['Demonstrates senior system design and full-stack implementation', 'Covers event-driven architecture, RAG, agentic workflows, memory fabric, governance, and observability', 'Provides a polished operator console and runnable demo', 'Strong FAANG-style architecture and portfolio case study']),
]
for idx, (title, subtitle, bullets) in enumerate(slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    if idx == 0:
        add_title(slide, title, subtitle)
        add_card(slide, 0.75, 1.65, 3.7, 1.5, 'Core Product Loop', 'Detect -> Normalize -> Correlate -> Retrieve -> Reason -> Guard -> Alert -> Explain', CYAN)
        add_card(slide, 4.85, 1.65, 3.7, 1.5, 'Primary Domain', 'Industrial safety, maintenance, security operations, and smart facilities', VIOLET)
        add_card(slide, 8.95, 1.65, 3.4, 1.5, 'Project Type', 'Senior-level full-stack AI architecture MVP', GREEN)
        add_bullets(slide, bullets, y=3.75, font=18)
    elif title == 'High-Level Architecture':
        add_title(slide, title, subtitle)
        cols = [('Sources','Cameras\nSensors\nAudio\nDocs'),('Ingest','FastAPI\nConnectors'),('Bus','Redpanda\nTopics'),('Process','Validate\nCorrelate'),('Memory','Postgres\nQdrant\nRedis'),('Agents','Investigate\nPlan\nGuard'),('UI','Console\nAudit')]
        x=0.45
        for name, body in cols:
            add_card(slide, x, 2.0, 1.65, 2.05, name, body, CYAN)
            x += 1.82
    else:
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets)
    # footer
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.25))
    p = foot.text_frame.paragraphs[0]; p.text = f'AMIF v1.0 - Complete Project Presentation - {idx+1}/{len(slides)}'; p.font.size = Pt(8); p.font.color.rgb = MUTED; p.alignment = PP_ALIGN.RIGHT

prs.save(ppt_path)

print('Generated:')
for f in ['AMIF_Complete_Project_Report.md','AMIF_Complete_Project_Report.html','AMIF_Complete_Project_Report.pdf','AMIF_Project_Presentation.pptx']:
    print(OUT / f)
